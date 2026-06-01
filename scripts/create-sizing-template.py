"""create-sizing-template.py - Generate sizing-base-template.pptx with baked donors.

Loads the Snowflake master PPTX template and produces a base presentation that
contains EIGHT fully-designed "donor" slides (cloned verbatim from the master),
keeping all slide masters, layouts, and themes intact, then BAKES default sizing
scaffolding into each so the raw template already reads as a sizing proposal:

    [0] title                (master idx 0)  - customer-first blue cover
    [1] agenda               (master idx 1)  - section list
    [2] safe_harbor          (master idx 2)  - Snowflake legal disclaimer (verbatim)
    [3] four_column_numbers  (master idx 11) - exec-summary big-number columns
    [4] table_styled         (master idx 18) - blue-header styled table
    [5] content              (master idx 7)  - one-column title + body
    [6] two_column           (master idx 20) - two paragraph-body columns
    [7] thank_you            (master idx 22) - blue closer w/ wordmark

The bake order matches renderer.pptx.clone.BAKED_DONOR_ORDER exactly, so at
render time build_pptx.py finds each donor by its slide index (not by sample
text), duplicates it in-package, re-injects real content (clone.py + inject.py),
and deletes the originals.  The baked defaults are placeholder text the renderer
overwrites; they exist so the committed .pptx reads correctly on its own.

Injection targets shapes by placeholder ROLE + POSITION (inject.py), so re-baking
the wording here does not break the renderer.

Donor indices use the master's SlideIdList ordering (the canonical deck order),
NOT the alphabetical slide-part ordering.  See snowflake-pptx TemplateMappings.

Usage:
    cd plugins/snowflake-sizing
    python3 scripts/create-sizing-template.py

Run once and commit assets/templates/sizing-base-template.pptx to version control.
"""
from __future__ import annotations

import pathlib
import sys

from pptx import Presentation

_SCRIPT_DIR = pathlib.Path(__file__).resolve().parent
_PLUGIN_ROOT = _SCRIPT_DIR.parent
_PPTX_PLUGIN = _PLUGIN_ROOT.parent / "snowflake-pptx"

# Make `renderer` importable when run from any CWD.
if str(_PLUGIN_ROOT) not in sys.path:
    sys.path.insert(0, str(_PLUGIN_ROOT))

from renderer.pptx import inject  # noqa: E402
from renderer.pptx.clone import (  # noqa: E402
    BAKED_DONOR_ORDER,
    clone_slide_crossfile,
    delete_slides,
)

SOURCE_TEMPLATE = _PPTX_PLUGIN / "assets" / "templates" / "SNOWFLAKE TEMPLATE JANUARY 2026.pptx"
OUT_PATH = _PLUGIN_ROOT / "assets" / "templates" / "sizing-base-template.pptx"

# Donor kind -> master SlideIdList index.  Iterated in BAKED_DONOR_ORDER so the
# output slide order stays in lock-step with the renderer's index-based lookup.
SRC_INDEX: dict[str, int] = {
    "title": 0,
    "agenda": 1,
    "safe_harbor": 2,
    "four_column_numbers": 11,
    "table_styled": 18,
    "content": 7,
    "two_column": 20,
    "thank_you": 22,
}

# Default scaffolding text baked into the donors (placeholder tokens the renderer
# overwrites).  Square-bracket tokens read as "fill me in" on the raw template.
_AGENDA_ITEMS = [
    "Executive Summary",
    "Warehouse Workloads",
    "Year-by-Year Costs",
    "Cost Detail by Year",
    "Serverless & AI / Cortex",
]
_WORKLOAD_HEADERS = [
    "Workload", "Size", "Hrs/Day", "Days/Mo", "Min Cl", "Max Cl", "Ramp", "Dev", "Go-Live",
]
_WORKLOAD_RATIOS = [2.5, 0.8, 0.85, 0.85, 0.75, 0.75, 1.05, 0.7, 0.8]
_WORKLOAD_SAMPLE = ["Analytics", "M", "8", "22", "1", "1", "linear", "M1", "M3"]
_EXEC_NUMBERS = ["$0.0M", "3 Yr", "$0.0M", "$0.00"]
_EXEC_CAPTIONS = ["Core TCV", "Contract Term", "Avg Annual Core", "Per-Credit Rate"]
_CONTENT_BODY = [
    "Pricing reflects current Snowflake list rates.",
    "Credit consumption estimated from workload sizing inputs.",
    "Storage assumes compression and the stated annual growth.",
    "Figures are planning estimates, not a commercial quote.",
]
_SERVERLESS_LEFT = [
    "Serverless Features",
    "- Snowpipe / Snowpipe Streaming",
    "- Materialized Views",
    "- Search Optimization",
    "",
    "Total (all years): $0",
]
_SERVERLESS_RIGHT = [
    "AI / Cortex Features",
    "- Cortex Analyst",
    "- Cortex Search",
    "- Cortex Complete",
    "",
    "Total (all years): $0",
]

# Donor kinds that carry a Confidential footer (content slides; the cover, safe
# harbor (own copyright) and closer are excluded).
_FOOTER_KINDS = {"agenda", "four_column_numbers", "table_styled", "content", "two_column"}


# ── Bake helpers ────────────────────────────────────────────────────────────── #

def _is_placeholder(shape) -> bool:
    try:
        pf = shape.placeholder_format
        return pf is not None and pf.type is not None
    except (ValueError, AttributeError):
        return False


def _bake_title(slide) -> None:
    # Drop the donor cover's right-side sample decorations / presenter block so
    # the customer-first title owns the slide (mirrors build_title_slide).
    inject.remove_shapes(slide, lambda sh: (sh.left or 0) > 5_500_000)
    shapes = inject.text_shapes_by_position(slide)  # main, subtitle, date (by top)
    if len(shapes) >= 1:
        inject.set_paragraph_texts(shapes[0], ["[Customer Name]", "Snowflake Sizing Proposal"])
    if len(shapes) >= 2:
        inject.set_shape_text(shapes[1], "[Edition]  |  [Cloud]  |  [Region]  |  [N]-Year Contract")
    if len(shapes) >= 3:
        inject.set_shape_text(shapes[2], "[Month Year]")


def _bake_agenda(slide) -> None:
    # Remove the donor's instructional sidebar (non-placeholder text auto-shape).
    inject.remove_shapes(slide, lambda sh: sh.has_text_frame and not _is_placeholder(sh))
    inject.set_title(slide, "Agenda")
    body = inject.find_largest_body(slide)
    if body is not None:
        inject.set_body_paragraphs(body, _AGENDA_ITEMS, tight=True)


def _bake_safe_harbor(slide) -> None:
    # The donor already carries the verbatim Snowflake safe-harbor legal text and
    # its own copyright line; keep it untouched.
    pass


def _bake_four_column(slide) -> None:
    inject.set_title(slide, "Executive Summary")
    inject.set_subtitle(slide, "[Edition]  |  [Cloud]  |  [Region]")
    for sh, val in zip(inject.number_shapes(slide), _EXEC_NUMBERS):
        inject.set_runs_font_size(sh, 2800, word_wrap=False)
        inject.set_shape_text(sh, val)
    for cap, label in zip(inject.caption_shapes(slide), _EXEC_CAPTIONS):
        inject.replace_caption(cap, label)


def _bake_table_styled(slide) -> None:
    inject.set_title(slide, "Warehouse Workloads")
    inject.set_subtitle(slide, "[N] workload(s)")
    inject.fill_table(slide, _WORKLOAD_HEADERS, [_WORKLOAD_SAMPLE], col_ratios=_WORKLOAD_RATIOS,
                      data_row_fill="FFFFFF")


def _bake_content(slide) -> None:
    inject.set_title(slide, "Key Assumptions")
    inject.set_subtitle(slide, "Pricing & sizing basis")
    bodies = inject.body_placeholders(slide)
    if bodies:
        inject.set_body_paragraphs(bodies[0], _CONTENT_BODY, tight=True)


def _bake_two_column(slide) -> None:
    inject.set_title(slide, "Serverless & AI / Cortex")
    inject.set_subtitle(slide, "Enabled features and projected spend")
    bodies = inject.body_placeholders(slide)
    if len(bodies) >= 2:
        inject.set_body_paragraphs(bodies[0], _SERVERLESS_LEFT, font_size=1300, tight=True)
        inject.set_body_paragraphs(bodies[1], _SERVERLESS_RIGHT, font_size=1300, tight=True)


def _bake_thank_you(slide) -> None:
    # Keep the iconic "THANK / YOU" headline; the renderer overlays a partnering
    # line + next steps at build time.
    pass


_BAKERS = {
    "title": _bake_title,
    "agenda": _bake_agenda,
    "safe_harbor": _bake_safe_harbor,
    "four_column_numbers": _bake_four_column,
    "table_styled": _bake_table_styled,
    "content": _bake_content,
    "two_column": _bake_two_column,
    "thank_you": _bake_thank_you,
}


def main() -> None:
    if not SOURCE_TEMPLATE.is_file():
        print(f"ERROR: source template not found:\n  {SOURCE_TEMPLATE}", file=sys.stderr)
        print("\nMake sure the snowflake-pptx plugin is present at the expected path.", file=sys.stderr)
        sys.exit(1)

    print(f"Loading source:  {SOURCE_TEMPLATE}")
    # `src` supplies donor slides to read from; `dst` is the output package.
    # Both load the full master so dst already owns every layout + theme; the
    # cross-file clone copies only the slide-level media it needs.
    src = Presentation(str(SOURCE_TEMPLATE))
    dst = Presentation(str(SOURCE_TEMPLATE))

    print(f"  Slide masters:   {len(dst.slide_masters)}")
    removed = delete_slides(dst, lambda _s: True)
    print(f"  Stripped {removed} designed slides (masters/layouts/themes kept)")

    print(f"\nBaking {len(BAKED_DONOR_ORDER)} donor slides (in render lookup order):")
    for kind in BAKED_DONOR_ORDER:
        idx = SRC_INDEX[kind]
        slide = clone_slide_crossfile(src, idx, dst)
        _BAKERS[kind](slide)
        if kind in _FOOTER_KINDS:
            inject.add_footer(slide, "Confidential")
        has_tbl = any(getattr(sh, "has_table", False) for sh in slide.shapes)
        title = inject.find_title_placeholder(slide)
        title_txt = title.text_frame.text.replace("\n", " / ")[:32] if title is not None else ""
        print(
            f"  + [{kind:>20}] master_idx={idx:<2} layout={slide.slide_layout.name!r:<26} "
            f"table={str(has_tbl):<5} title={title_txt!r}"
        )

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    dst.save(str(OUT_PATH))
    size_kb = OUT_PATH.stat().st_size // 1024
    print(f"\nSaved: {OUT_PATH}  ({size_kb} KB, {len(list(dst.slides))} donor slides)")
    print("Done. Commit assets/templates/sizing-base-template.pptx to version control.")


if __name__ == "__main__":
    main()
